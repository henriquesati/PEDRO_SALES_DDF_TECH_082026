from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pathlib import Path
from PIL import Image, ImageDraw, ImageFont

BASE_DIR = Path(__file__).resolve().parent
ROOT_DIR = BASE_DIR.parent
OUT = BASE_DIR / "arquitetura_dadosfera.pptx"
ASSETS_DIR = ROOT_DIR / "assets"
ICONS_DIR = ASSETS_DIR / "icons"
ASSETS_DIR.mkdir(parents=True, exist_ok=True)
ICONS_DIR.mkdir(parents=True, exist_ok=True)
W, H = 13.333, 7.5
C = {
    "navy": RGBColor(30, 58, 138), "blue": RGBColor(37, 99, 235),
    "indigo": RGBColor(79, 70, 229), "green": RGBColor(5, 150, 105),
    "red": RGBColor(229, 62, 62), "orange": RGBColor(221, 107, 32),
    "ink": RGBColor(45, 55, 72), "muted": RGBColor(100, 116, 139),
    "line": RGBColor(203, 213, 225), "paper": RGBColor(248, 250, 252),
    "white": RGBColor(255, 255, 255), "softred": RGBColor(254, 242, 242),
    "softblue": RGBColor(239, 246, 255), "softgreen": RGBColor(236, 253, 245),
}

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(W), Inches(H)
blank = prs.slide_layouts[6]


def create_assets():
    """Gera PNGs individuais, leves e transparentes para animar no PowerPoint."""
    names = {
        "kinesis": "K", "firehose": "F", "lambda": "λ", "sqs": "Q", "dynamodb": "D", "s3": "S3",
        "glue": "G", "redshift": "R", "airflow": "A", "redis": "R", "docker": "D", "terraform": "T",
        "github-actions": "GH", "secrets-manager": "SM", "lake-formation": "LF", "datahub": "DH",
        "powerbi": "PBI", "tableau": "T", "datadog": "DD", "eventbridge": "EB", "maestro-api": "M",
        "snowflake": "SF", "metabase": "MB", "streamlit": "ST", "genai": "AI", "dadosfera": "DF",
    }
    font_paths = ["/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", "/usr/share/fonts/dejavu/DejaVuSans-Bold.ttf"]
    font_path = next((p for p in font_paths if Path(p).exists()), None)
    for name, label in names.items():
        image = Image.new("RGBA", (256, 256), (255, 255, 255, 0))
        draw = ImageDraw.Draw(image)
        draw.rounded_rectangle((18, 18, 238, 238), radius=42, fill=(30, 58, 138, 255), outline=(255, 255, 255, 255), width=5)
        font = ImageFont.truetype(font_path, 64) if font_path else ImageFont.load_default()
        bounds = draw.textbbox((0, 0), label, font=font)
        draw.text(((256-(bounds[2]-bounds[0]))/2, (256-(bounds[3]-bounds[1]))/2-8), label, font=font, fill=(255, 255, 255, 255))
        image.save(ICONS_DIR / f"{name}.png")


def create_chart_png(filename, title, subtitle, titles, content, accent, panel_fill, pains=None, solution=False):
    """Exporta o gráfico L2R como PNG independente, sem depender dos ícones."""
    scale, width, height = 2, 1600, 900
    image = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(image)
    font_candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
        "/usr/share/fonts/dejavu/DejaVuSans.ttf",
    ]
    bold_candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
        "/usr/share/fonts/truetype/liberation2/LiberationSans-Bold.ttf",
        "/usr/share/fonts/dejavu/DejaVuSans-Bold.ttf",
    ]
    font_path = next((p for p in font_candidates if Path(p).exists()), None)
    bold_path = next((p for p in bold_candidates if Path(p).exists()), font_path)
    def f(size, bold=False):
        path = bold_path if bold else font_path
        return ImageFont.truetype(path, size) if path else ImageFont.load_default()
    draw.text((66, 38), title, font=f(36, True), fill=(45, 55, 72))
    draw.text((66, 88), subtitle, font=f(18), fill=(100, 116, 139))
    draw.rectangle((66, 132, 1534, 138), fill=accent)
    top, left, gap, cw, ch = 175, 66, 10, 308, 570
    for i, heading in enumerate(titles):
        x = left + i * (cw + gap)
        draw.rounded_rectangle((x, top, x+cw, top+ch), radius=12, fill=panel_fill, outline=(203, 213, 225), width=2)
        draw.rounded_rectangle((x, top, x+cw, top+74), radius=12, fill=accent)
        draw.rectangle((x, top+55, x+cw, top+74), fill=accent)
        draw.text((x+18, top+18), str(i+1), font=f(20, True), fill="white")
        draw.multiline_text((x+58, top+13), heading, font=f(18, True), fill="white", spacing=2)
        y = top + 98
        for label in content[i]:
            draw.rounded_rectangle((x+18, y, x+cw-18, y+48), radius=8, fill="white", outline=(203, 213, 225), width=2)
            draw.text((x+30, y+12), label, font=f(15, True), fill=(45, 55, 72))
            y += 59
        if pains:
            for pain in pains[i]:
                draw.rounded_rectangle((x+18, y+6, x+cw-18, y+42), radius=8, fill=(254, 242, 242), outline=(229, 62, 62), width=2)
                draw.text((x+30, y+14), pain, font=f(13, True), fill=(229, 62, 62))
                y += 46
        if i < 4:
            mid = top + 37
            draw.line((x+cw+2, mid, x+cw+gap-2, mid), fill=accent, width=4)
            draw.polygon([(x+cw+gap-2, mid), (x+cw+gap-12, mid-7), (x+cw+gap-12, mid+7)], fill=accent)
    image.save(ASSETS_DIR / filename, optimize=True)


def box(slide, x, y, w, h, fill, line=None, radius=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(0.8)
    return shape


def text(slide, value, x, y, w, h, size=10, color=None, bold=False, align=PP_ALIGN.LEFT, font="Aptos", valign=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True; tf.margin_left = Inches(0.07); tf.margin_right = Inches(0.07); tf.margin_top = Inches(0.03); tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = value; r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color or C["ink"]
    return tb


def pill(slide, label, x, y, w, color, fill):
    box(slide, x, y, w, 0.28, fill, fill)
    text(slide, label, x, y+0.01, w, 0.23, 7.2, color, True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def footer(slide, label, color):
    text(slide, label, 0.55, 7.16, 12.2, 0.18, 7, color, True, PP_ALIGN.RIGHT)


def header(slide, kicker, title, subtitle, accent):
    text(slide, kicker.upper(), 0.55, 0.25, 5.5, 0.25, 8, accent, True)
    text(slide, title, 0.55, 0.52, 12.2, 0.48, 24, C["ink"], True)
    text(slide, subtitle, 0.55, 1.05, 12.2, 0.3, 9.5, C["muted"])
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.42), Inches(12.23), Inches(0.025)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = accent; slide.shapes[-1].line.fill.background()


def columns(slide, titles, accent, content, pains=None, solution=False):
    left, top, gap, cw, ch = 0.55, 1.7, 0.08, 2.38, 4.72
    for i, title_ in enumerate(titles):
        x = left + i*(cw+gap)
        fill = C["softblue"] if solution else C["paper"]
        box(slide, x, top, cw, ch, fill, C["line"])
        box(slide, x, top, cw, 0.62, accent if solution else C["ink"], accent if solution else C["ink"])
        text(slide, f"{i+1}", x+0.12, top+0.1, 0.25, 0.25, 10, C["white"], True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        text(slide, title_, x+0.4, top+0.09, cw-0.5, 0.46, 9, C["white"], True, PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE)
        y = top + 0.78
        for label in content[i]:
            box(slide, x+0.13, y, cw-0.26, 0.42, C["white"], C["line"])
            text(slide, label, x+0.2, y+0.04, cw-0.4, 0.33, 8.2, C["ink"], True)
            y += 0.49
        if pains:
            for p in pains[i]:
                pill(slide, p, x+0.13, y+0.06, cw-0.26, C["red"], C["softred"])
                y += 0.35
        if i < 4:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x+cw), Inches(top+0.36), Inches(x+cw+gap), Inches(top+0.36))
            line.line.color.rgb = accent; line.line.width = Pt(1.4)


def add_legacy():
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C["white"]
    header(slide, "Ato 1 · diagnóstico", "Arquitetura Legada — Complexidade & Risco", "AWS DIY: muitas ferramentas, muitas colas manuais, pouca previsibilidade operacional.", C["red"])
    titles = ["Ingestão &\nValidação", "Processamento,\nLimpeza & DW", "Orquestração &\nInfra / DevOps", "Catálogo,\nMetadados & LGPD", "Consumo,\nAnalytics & Alertas"]
    content = [
        ["Kinesis Streams", "Kinesis Firehose", "Lambda + Great Expectations", "SQS DLQ", "DynamoDB / S3 Quarentena", "CloudWatch + SNS"],
        ["S3 Bronze", "Glue Jobs / PySpark", "S3 Silver", "Redshift Cluster", "S3 Gold"],
        ["MWAA / Airflow", "ElastiCache Redis", "Docker + ECR", "Terraform / CloudFormation", "GitHub Actions / GitLab CI", "Secrets Manager"],
        ["Glue Crawlers", "Lake Formation", "IAM Policies JSON", "DataHub / OpenMetadata"],
        ["PowerBI / Tableau", "Datadog / Grafana", "EventBridge + SNS", "Custom Lambda / Python APIs"],
    ]
    pains = [
        ["Sharding manual", "Schema frágil"], ["Cold start 1–4 min", "DPUs ociosos"], ["Downtime no pico", "+1 Platform Engineer"], ["IAM complexo", "Risco LGPD"], ["3–6 semanas", "Foco em manutenção"],
    ]
    columns(slide, titles, C["red"], content, pains)
    box(slide, 0.55, 6.58, 12.23, 0.45, C["softred"], C["red"])
    text(slide, "GARGALO", 0.72, 6.67, 0.85, 0.18, 8, C["red"], True)
    text(slide, "Lead time: 3–6 semanas   ·   Sustentação: 1 Platform Engineer + 1–2 Data Engineers   ·   Impacto: engenharia presa em servidores e pipelines quebrados", 1.55, 6.65, 10.9, 0.22, 8.4, C["ink"], True)
    create_chart_png("grafico-legado-l2r.png", "Arquitetura Legada — Complexidade & Risco", "AWS DIY: muitas ferramentas, muitas colas manuais, pouca previsibilidade operacional.", titles, content, (229, 62, 62), (248, 250, 252), pains)
    footer(slide, "01 / 02", C["red"])


def add_solution():
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C["white"]
    header(slide, "Ato 1 · virada", "Arquitetura Proposta — Plataforma Dadosfera", "Um Sistema Operacional de Dados: absorve as camadas e devolve velocidade, governança e foco no negócio.", C["blue"])
    box(slide, 0.55, 1.56, 12.23, 0.32, C["navy"], C["navy"])
    text(slide, "PLATAFORMA DADOSFERA  ·  SaaS ALL-IN-ONE / SNOWFLAKE LAKEHOUSE", 0.75, 1.6, 11.8, 0.22, 9, C["white"], True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    titles = ["Ingestão &\nQualidade", "Processamento\n& DW", "Orquestração &\nInfraestrutura", "Catálogo, Metadados\n& Governança", "Consumo, Analytics\n& GenAI"]
    content = [
        ["API Maestro", "Great Expectations nativo", "Quarentena automática", "Schema validation"],
        ["Snowflake DW", "Star Schema Kimball", "Silver conformada", "Visões Gold"],
        ["SaaS gerenciado", "Sem servidores", "Sem sharding", "Elasticidade em picos"],
        ["Catálogo nativo", "Data Asset IDs", "RBAC + LGPD", "Rastreabilidade ponta a ponta"],
        ["Metabase integrado", "Streamlit Data Apps", "GenAI contextual", "Alertas self-service"],
    ]
    columns(slide, titles, C["blue"], content, solution=True)
    box(slide, 0.55, 6.58, 12.23, 0.45, C["softgreen"], C["green"])
    text(slide, "GANHOS", 0.72, 6.67, 0.85, 0.18, 8, C["green"], True)
    text(slide, "Lead time: < 3 dias (-86%)   ·   Zero risco de sharding / downtime em picos   ·   Equipe focada em receita e ROI   ·   SaaS previsível", 1.55, 6.65, 10.9, 0.22, 8.4, C["ink"], True)
    create_chart_png("grafico-dadosfera-l2r.png", "Arquitetura Proposta — Plataforma Dadosfera", "Um Sistema Operacional de Dados: absorve as camadas e devolve velocidade, governança e foco no negócio.", titles, content, (37, 99, 235), (239, 246, 255), solution=True)
    footer(slide, "02 / 02", C["blue"])

create_assets()
add_legacy(); add_solution()
prs.save(OUT)
print(f"Arquivo gerado: {OUT}")
print(f"Slides: {len(prs.slides)}")
